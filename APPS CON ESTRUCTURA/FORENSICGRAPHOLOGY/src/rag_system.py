import os
import re
import fitz  # PyMuPDF
import docx
import pptx
import numpy as np
import pandas as pd
import chromadb
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import Chroma
from langchain.schema import Document
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.llms import HuggingFaceHub
from sentence_transformers import SentenceTransformer
import torch
import re
import hashlib
import json
import datetime


class DocumentProcessor:
    """
    Classe per l'elaborazione e l'estrazione di testo da vari formati di documenti.
    """
    
    def __init__(self, upload_dir):
        """
        Inizializza il processore di documenti.
        
        Args:
            upload_dir (str): Directory dove salvare i documenti caricati
        """
        self.upload_dir = upload_dir
        os.makedirs(upload_dir, exist_ok=True)
    
    def save_uploaded_file(self, file_obj, filename=None):
        """
        Salva un file caricato nella directory di upload.
        
        Args:
            file_obj: Oggetto file caricato
            filename (str, optional): Nome del file
            
        Returns:
            str: Percorso del file salvato
        """
        if filename is None:
            filename = file_obj.name
        
        # Genera un nome file sicuro
        safe_filename = self._sanitize_filename(filename)
        
        # Aggiungi timestamp per evitare sovrascritture
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        filename_with_timestamp = f"{timestamp}_{safe_filename}"
        
        # Percorso completo del file
        file_path = os.path.join(self.upload_dir, filename_with_timestamp)
        
        # Salva il file
        with open(file_path, 'wb') as f:
            f.write(file_obj.read())
        
        return file_path
    
    def _sanitize_filename(self, filename):
        """
        Sanitizza un nome file rimuovendo caratteri non sicuri.
        
        Args:
            filename (str): Nome file originale
            
        Returns:
            str: Nome file sanitizzato
        """
        # Rimuovi caratteri non sicuri
        safe_filename = re.sub(r'[^\w\.-]', '_', filename)
        return safe_filename
    
    def extract_text(self, file_path):
        """
        Estrae il testo da un file in base al suo formato.
        
        Args:
            file_path (str): Percorso del file
            
        Returns:
            str: Testo estratto
        """
        # Determina il formato del file dall'estensione
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif ext == '.txt':
            return self.extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Formato file non supportato: {ext}")
    
    def extract_text_from_pdf(self, pdf_path):
        """
        Estrae il testo da un file PDF.
        
        Args:
            pdf_path (str): Percorso del file PDF
            
        Returns:
            str: Testo estratto
        """
        text = ""
        try:
            # Apri il documento PDF
            doc = fitz.open(pdf_path)
            
            # Estrai il testo da ogni pagina
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text += page.get_text()
            
            # Chiudi il documento
            doc.close()
        except Exception as e:
            print(f"Errore nell'estrazione del testo dal PDF {pdf_path}: {e}")
        
        return text
    
    def extract_text_from_docx(self, docx_path):
        """
        Estrae il testo da un file DOCX.
        
        Args:
            docx_path (str): Percorso del file DOCX
            
        Returns:
            str: Testo estratto
        """
        text = ""
        try:
            # Apri il documento DOCX
            doc = docx.Document(docx_path)
            
            # Estrai il testo da ogni paragrafo
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Estrai il testo dalle tabelle
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
        except Exception as e:
            print(f"Errore nell'estrazione del testo dal DOCX {docx_path}: {e}")
        
        return text
    
    def extract_text_from_pptx(self, pptx_path):
        """
        Estrae il testo da un file PPTX.
        
        Args:
            pptx_path (str): Percorso del file PPTX
            
        Returns:
            str: Testo estratto
        """
        text = ""
        try:
            # Apri la presentazione PPTX
            prs = pptx.Presentation(pptx_path)
            
            # Estrai il testo da ogni diapositiva
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Errore nell'estrazione del testo dal PPTX {pptx_path}: {e}")
        
        return text
    
    def extract_text_from_txt(self, txt_path):
        """
        Estrae il testo da un file TXT.
        
        Args:
            txt_path (str): Percorso del file TXT
            
        Returns:
            str: Testo estratto
        """
        try:
            # Apri il file TXT
            with open(txt_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            # Prova con una codifica diversa
            try:
                with open(txt_path, 'r', encoding='latin-1') as f:
                    text = f.read()
            except Exception as e:
                print(f"Errore nell'estrazione del testo dal TXT {txt_path}: {e}")
                text = ""
        except Exception as e:
            print(f"Errore nell'estrazione del testo dal TXT {txt_path}: {e}")
            text = ""
        
        return text
    
    def chunk_text(self, text, chunk_size=500, chunk_overlap=50):
        """
        Divide il testo in chunk più piccoli.
        
        Args:
            text (str): Testo da dividere
            chunk_size (int): Dimensione di ogni chunk in token
            chunk_overlap (int): Sovrapposizione tra chunk consecutivi
            
        Returns:
            list: Lista di chunk di testo
        """
        # Utilizza il text splitter di LangChain
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len
        )
        
        # Dividi il testo in chunk
        chunks = text_splitter.split_text(text)
        
        return chunks
    
    def process_document(self, file_path, chunk_size=500, chunk_overlap=50):
        """
        Elabora un documento: estrae il testo e lo divide in chunk.
        
        Args:
            file_path (str): Percorso del file
            chunk_size (int): Dimensione di ogni chunk in token
            chunk_overlap (int): Sovrapposizione tra chunk consecutivi
            
        Returns:
            dict: Informazioni sul documento elaborato
        """
        # Estrai il testo dal documento
        text = self.extract_text(file_path)
        
        # Dividi il testo in chunk
        chunks = self.chunk_text(text, chunk_size, chunk_overlap)
        
        # Calcola l'hash del file per l'identificazione
        file_hash = self._calculate_file_hash(file_path)
        
        # Ottieni il nome del file
        filename = os.path.basename(file_path)
        
        # Crea metadati per il documento
        metadata = {
            'filename': filename,
            'file_path': file_path,
            'file_hash': file_hash,
            'chunk_count': len(chunks),
            'total_text_length': len(text),
            'processing_date': datetime.datetime.now().isoformat()
        }
        
        return {
            'text': text,
            'chunks': chunks,
            'metadata': metadata
        }
    
    def _calculate_file_hash(self, file_path):
        """
        Calcola l'hash SHA-256 di un file.
        
        Args:
            file_path (str): Percorso del file
            
        Returns:
            str: Hash SHA-256 del file
        """
        sha256_hash = hashlib.sha256()
        
        with open(file_path, "rb") as f:
            # Leggi il file a blocchi per gestire file di grandi dimensioni
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        
        return sha256_hash.hexdigest()


class VectorStore:
    """
    Classe per la gestione del vector store per il sistema RAG.
    """
    
    def __init__(self, persist_directory, embedding_model_name="all-MiniLM-L6-v2"):
        """
        Inizializza il vector store.
        
        Args:
            persist_directory (str): Directory dove salvare il vector store
            embedding_model_name (str): Nome del modello di embedding
        """
        self.persist_directory = persist_directory
        self.embedding_model_name = embedding_model_name
        
        # Crea la directory se non esiste
        os.makedirs(persist_directory, exist_ok=True)
        
        # Inizializza il modello di embedding
        self.embedding_model = self._initialize_embedding_model(embedding_model_name)
        
        # Inizializza il vector store
        self.vector_store = self._initialize_vector_store()
    
    def _initialize_embedding_model(self, model_name):
        """
        Inizializza il modello di embedding.
        
        Args:
            model_name (str): Nome del modello
            
        Returns:
            object: Modello di embedding
        """
        try:
            # Utilizza HuggingFaceEmbeddings di LangChain
            embedding_model = HuggingFaceEmbeddings(model_name=model_name)
            return embedding_model
        except Exception as e:
            print(f"Errore nell'inizializzazione del modello di embedding: {e}")
            # Fallback: carica direttamente il modello con sentence-transformers
            return SentenceTransformer(model_name)
    
    def _initialize_vector_store(self):
        """
        Inizializza il vector store.
        
        Returns:
            object: Vector store
        """
        try:
            # Controlla se esiste già un vector store
            if os.path.exists(os.path.join(self.persist_directory, 'chroma.sqlite3')):
                # Carica il vector store esistente
                vector_store = Chroma(
                    persist_directory=self.persist_directory,
                    embedding_function=self.embedding_model
                )
            else:
                # Crea un nuovo vector store
                vector_store = Chroma(
                    persist_directory=self.persist_directory,
                    embedding_function=self.embedding_model
                )
            
            return vector_store
        except Exception as e:
            print(f"Errore nell'inizializzazione del vector store: {e}")
            # Fallback: utilizza direttamente ChromaDB
            client = chromadb.PersistentClient(path=self.persist_directory)
            collection_name = "forensic_graphology_docs"
            
            # Controlla se la collezione esiste già
            try:
                collection = client.get_collection(name=collection_name)
            except:
                # Crea una nuova collezione
                collection = client.create_collection(name=collection_name)
            
            return collection
    
    def add_document(self, document_info):
        """
        Aggiunge un documento al vector store.
        
        Args:
            document_info (dict): Informazioni sul documento
            
        Returns:
            dict: Risultato dell'operazione
        """
        chunks = document_info['chunks']
        metadata = document_info['metadata']
        
        # Crea documenti LangChain
        documents = []
        for i, chunk in enumerate(chunks):
            # Crea metadati per il chunk
            chunk_metadata = metadata.copy()
            chunk_metadata['chunk_id'] = i
            chunk_metadata['chunk_index'] = i
            chunk_metadata['chunk_total'] = len(chunks)
            
            # Crea un documento LangChain
            doc = Document(page_content=chunk, metadata=chunk_metadata)
            documents.append(doc)
        
        try:
            # Aggiungi i documenti al vector store
            self.vector_store.add_documents(documents)
            
            return {
                'success': True,
                'document_id': metadata['file_hash'],
                'chunks_added': len(chunks)
            }
        except Exception as e:
            print(f"Errore nell'aggiunta del documento al vector store: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def search(self, query, k=4):
        """
        Cerca documenti simili alla query.
        
        Args:
            query (str): Query di ricerca
            k (int): Numero di risultati da restituire
            
        Returns:
            list: Lista di documenti simili
        """
        try:
            # Cerca documenti simili
            results = self.vector_store.similarity_search(query, k=k)
            return results
        except Exception as e:
            print(f"Errore nella ricerca: {e}")
            return []
    
    def delete_document(self, document_id):
        """
        Elimina un documento dal vector store.
        
        Args:
            document_id (str): ID del documento
            
        Returns:
            dict: Risultato dell'operazione
        """
        try:
            # Elimina il documento
            self.vector_store.delete(filter={"file_hash": document_id})
            
            return {
                'success': True,
                'document_id': document_id
            }
        except Exception as e:
            print(f"Errore nell'eliminazione del documento: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def get_all_documents(self):
        """
        Ottiene tutti i documenti nel vector store.
        
        Returns:
            list: Lista di documenti
        """
        try:
            # Ottieni tutti i documenti
            results = self.vector_store.get()
            
            # Estrai i metadati unici
            unique_docs = {}
            for i, metadata in enumerate(results['metadatas']):
                if 'file_hash' in metadata:
                    file_hash = metadata['file_hash']
                    if file_hash not in unique_docs:
                        unique_docs[file_hash] = {
                            'document_id': file_hash,
                            'filename': metadata.get('filename', 'Unknown'),
                            'file_path': metadata.get('file_path', ''),
                            'chunk_total': metadata.get('chunk_total', 0),
                            'processing_date': metadata.get('processing_date', '')
                        }
            
            return list(unique_docs.values())
        except Exception as e:
            print(f"Errore nel recupero dei documenti: {e}")
            return []


class RAGSystem:
    """
    Classe per il sistema RAG (Retrieval Augmented Generation).
    """
    
    def __init__(self, upload_dir, vector_store_dir, use_local_model=False, model_name=None):
        """
        Inizializza il sistema RAG.
        
        Args:
            upload_dir (str): Directory per i documenti caricati
            vector_store_dir (str): Directory per il vector store
            use_local_model (bool): Se utilizzare un modello locale
            model_name (str): Nome del modello da utilizzare
        """
        self.document_processor = DocumentProcessor(upload_dir)
        self.vector_store = VectorStore(vector_store_dir)
        self.use_local_model = use_local_model
        self.model_name = model_name
        
        # Inizializza il modello come None (modalità senza LLM)
        self.model = None
        
        # Prova a inizializzare il modello solo se specificato
        if model_name:
            try:
                self._initialize_model(use_local_model, model_name)
            except Exception as e:
                print(f"Errore nell'inizializzazione del modello: {e}")
                print("Il sistema RAG funzionerà in modalità di sola ricerca (senza generazione).")
    
    def _initialize_model(self, use_local_model, model_name):
        """
        Inizializza il modello di linguaggio.
        
        Args:
            use_local_model (bool): Se utilizzare un modello locale
            model_name (str): Nome del modello
            
        Returns:
            object: Modello di linguaggio
        """
        # In questa versione semplificata, non inizializziamo alcun modello
        # per evitare problemi di dipendenze e token API
        print("Modalità di sola ricerca attivata (senza generazione).")
        return None
    
    def process_and_store_document(self, file_obj, filename=None):
        """
        Elabora e memorizza un documento.
        
        Args:
            file_obj: Oggetto file caricato
            filename (str, optional): Nome del file
            
        Returns:
            dict: Risultato dell'operazione
        """
        try:
            # Salva il file caricato
            file_path = self.document_processor.save_uploaded_file(file_obj, filename)
            
            # Elabora il documento
            document_info = self.document_processor.process_document(file_path)
            
            # Aggiungi il documento al vector store
            result = self.vector_store.add_document(document_info)
            
            # Aggiungi informazioni aggiuntive al risultato
            result['filename'] = os.path.basename(file_path)
            result['file_path'] = file_path
            result['chunk_count'] = len(document_info['chunks'])
            
            return result
        except Exception as e:
            print(f"Errore nell'elaborazione e memorizzazione del documento: {e}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def query(self, query_text, k=4, scrub_sensitive=True):
        """
        Esegue una query sul sistema RAG.
        
        Args:
            query_text (str): Testo della query
            k (int): Numero di documenti da recuperare
            scrub_sensitive (bool): Se rimuovere informazioni sensibili
            
        Returns:
            dict: Risultato della query
        """
        try:
            # Cerca documenti simili
            retrieved_docs = self.vector_store.search(query_text, k=k)
            
            # Estrai il contesto dai documenti
            context = self._build_context(retrieved_docs, scrub_sensitive)
            
            # Prepara i riferimenti
            references = self._prepare_references(retrieved_docs)
            
            # Se non c'è un modello, restituisci solo i documenti recuperati
            if self.model is None:
                response = "Modalità di sola ricerca attiva. Ecco i documenti più rilevanti per la tua query:\n\n"
                for i, doc in enumerate(retrieved_docs):
                    response += f"[Documento {i+1}] {doc.metadata.get('filename', 'Unknown')}\n"
                    response += f"Estratto: {doc.page_content[:200]}...\n\n"
            else:
                # Crea il prompt
                prompt = self._create_prompt(query_text, context)
                
                # Genera la risposta
                response = self._generate_response(prompt)
            
            return {
                'success': True,
                'query': query_text,
                'response': response,
                'references': references
            }
        except Exception as e:
            print(f"Errore nell'esecuzione della query: {e}")
            return {
                'success': False,
                'error': str(e),
                'query': query_text
            }
    
    def _build_context(self, documents, scrub_sensitive=True):
        """
        Costruisce il contesto dai documenti recuperati.
        
        Args:
            documents (list): Lista di documenti
            scrub_sensitive (bool): Se rimuovere informazioni sensibili
            
        Returns:
            str: Contesto
        """
        context_parts = []
        
        for i, doc in enumerate(documents):
            # Estrai il contenuto e i metadati
            content = doc.page_content
            metadata = doc.metadata
            
            # Rimuovi informazioni sensibili se richiesto
            if scrub_sensitive:
                content = self._scrub_sensitive_info(content)
            
            # Aggiungi il contenuto al contesto
            context_parts.append(f"[Documento {i+1}] {content}")
        
        # Unisci le parti del contesto
        context = "\n\n".join(context_parts)
        
        return context
    
    def _scrub_sensitive_info(self, text):
        """
        Rimuove informazioni sensibili dal testo.
        
        Args:
            text (str): Testo da elaborare
            
        Returns:
            str: Testo elaborato
        """
        # Rimuovi numeri di telefono
        text = re.sub(r'\b\d{10}\b', '[TELEFONO]', text)
        text = re.sub(r'\b\d{3}[-.\s]?\d{3}[-.\s]?\d{4}\b', '[TELEFONO]', text)
        
        # Rimuovi indirizzi email
        text = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[EMAIL]', text)
        
        # Rimuovi codici fiscali italiani
        text = re.sub(r'\b[A-Z]{6}\d{2}[A-Z]\d{2}[A-Z]\d{3}[A-Z]\b', '[CODICE_FISCALE]', text)
        
        # Rimuovi numeri di carte di credito
        text = re.sub(r'\b\d{4}[-\s]?\d{4}[-\s]?\d{4}[-\s]?\d{4}\b', '[CARTA_DI_CREDITO]', text)
        
        # Rimuovi IBAN
        text = re.sub(r'\b[A-Z]{2}\d{2}[A-Z0-9]{4}\d{7}[A-Z0-9]{0,16}\b', '[IBAN]', text)
        
        return text
    
    def _create_prompt(self, query, context):
        """
        Crea il prompt per il modello.
        
        Args:
            query (str): Query dell'utente
            context (str): Contesto dai documenti
            
        Returns:
            str: Prompt
        """
        prompt_template = """
        Sei un consulente di Grafologia Forense. Ti fornisco del contesto da documenti caricati.
        Rispondi in modo coerente e professionale, senza rivelare mai dati privati.
        
        CONTENUTO RILEVANTE:
        {context}
        
        DOMANDA: {query}
        
        RISPOSTA:
        """
        
        # Crea il prompt
        prompt = PromptTemplate(
            template=prompt_template,
            input_variables=["context", "query"]
        )
        
        # Formatta il prompt
        formatted_prompt = prompt.format(context=context, query=query)
        
        return formatted_prompt
    
    def _generate_response(self, prompt):
        """
        Genera una risposta dal modello.
        
        Args:
            prompt (str): Prompt per il modello
            
        Returns:
            str: Risposta generata
        """
        if self.model is None:
            return "Mi dispiace, il modello di linguaggio non è disponibile al momento."
        
        try:
            # Crea una chain
            chain = LLMChain(llm=self.model, prompt=PromptTemplate.from_template(prompt))
            
            # Genera la risposta
            response = chain.run("")
            
            return response
        except Exception as e:
            print(f"Errore nella generazione della risposta: {e}")
            
            # Fallback: risposta semplice
            return "Mi dispiace, non sono riuscito a generare una risposta. Si è verificato un errore."
    
    def _prepare_references(self, documents):
        """
        Prepara i riferimenti ai documenti.
        
        Args:
            documents (list): Lista di documenti
            
        Returns:
            list: Lista di riferimenti
        """
        references = []
        
        for i, doc in enumerate(documents):
            # Estrai i metadati
            metadata = doc.metadata
            
            # Crea un riferimento
            reference = {
                'id': i + 1,
                'filename': metadata.get('filename', 'Unknown'),
                'chunk_id': metadata.get('chunk_id', 0),
                'chunk_index': metadata.get('chunk_index', 0),
                'chunk_total': metadata.get('chunk_total', 0),
                'snippet': doc.page_content[:100] + "..." if len(doc.page_content) > 100 else doc.page_content
            }
            
            references.append(reference)
        
        return references
    
    def get_document_list(self):
        """
        Ottiene la lista dei documenti memorizzati.
        
        Returns:
            list: Lista di documenti
        """
        return self.vector_store.get_all_documents()
    
    def delete_document(self, document_id):
        """
        Elimina un documento.
        
        Args:
            document_id (str): ID del documento
            
        Returns:
            dict: Risultato dell'operazione
        """
        return self.vector_store.delete_document(document_id)
